"""
PowerPoint Generator for Creating Beautiful Presentations
Professional-grade slide generation with modern design elements.
"""
import logging
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

logger = logging.getLogger(__name__)


class PPTXGenerator:
    """Generator for creating polished PowerPoint presentations."""

    THEMES = {
        "professional": {
            "bg_color": RGBColor(255, 255, 255),
            "title_bg": RGBColor(19, 47, 76),
            "title_color": RGBColor(255, 255, 255),
            "slide_title_color": RGBColor(19, 47, 76),
            "text_color": RGBColor(55, 65, 81),
            "accent": RGBColor(37, 99, 235),
            "accent_light": RGBColor(219, 234, 254),
            "bullet_color": RGBColor(37, 99, 235),
            "subtitle_color": RGBColor(180, 198, 220),
            "font_heading": "Calibri",
            "font_body": "Calibri",
        },
        "modern": {
            "bg_color": RGBColor(250, 250, 250),
            "title_bg": RGBColor(17, 24, 39),
            "title_color": RGBColor(255, 255, 255),
            "slide_title_color": RGBColor(17, 24, 39),
            "text_color": RGBColor(75, 85, 99),
            "accent": RGBColor(139, 92, 246),
            "accent_light": RGBColor(237, 233, 254),
            "bullet_color": RGBColor(139, 92, 246),
            "subtitle_color": RGBColor(156, 163, 175),
            "font_heading": "Segoe UI",
            "font_body": "Segoe UI",
        },
        "minimal": {
            "bg_color": RGBColor(255, 255, 255),
            "title_bg": RGBColor(24, 24, 27),
            "title_color": RGBColor(255, 255, 255),
            "slide_title_color": RGBColor(24, 24, 27),
            "text_color": RGBColor(82, 82, 91),
            "accent": RGBColor(24, 24, 27),
            "accent_light": RGBColor(244, 244, 245),
            "bullet_color": RGBColor(161, 161, 170),
            "subtitle_color": RGBColor(161, 161, 170),
            "font_heading": "Calibri Light",
            "font_body": "Calibri",
        },
        "dark": {
            "bg_color": RGBColor(17, 24, 39),
            "title_bg": RGBColor(0, 0, 0),
            "title_color": RGBColor(0, 255, 136),
            "slide_title_color": RGBColor(0, 255, 136),
            "text_color": RGBColor(209, 213, 219),
            "accent": RGBColor(0, 255, 136),
            "accent_light": RGBColor(30, 41, 59),
            "bullet_color": RGBColor(0, 255, 136),
            "subtitle_color": RGBColor(107, 114, 128),
            "font_heading": "Consolas",
            "font_body": "Segoe UI",
        },
        "cyber": {
            "bg_color": RGBColor(13, 17, 23),
            "title_bg": RGBColor(0, 0, 0),
            "title_color": RGBColor(0, 212, 255),
            "slide_title_color": RGBColor(0, 212, 255),
            "text_color": RGBColor(201, 209, 217),
            "accent": RGBColor(0, 212, 255),
            "accent_light": RGBColor(22, 27, 34),
            "bullet_color": RGBColor(0, 212, 255),
            "subtitle_color": RGBColor(139, 148, 158),
            "font_heading": "Consolas",
            "font_body": "Segoe UI",
        },
    }

    def __init__(self, theme: str = "professional"):
        self.theme_name = theme
        self.theme = self.THEMES.get(theme, self.THEMES["professional"])
        logger.info(f"PPTXGenerator initialized with theme: {theme}")

    def create_presentation(self, title: str, slides_content: list, output_path: str) -> str:
        """Create a PowerPoint presentation and save it."""
        try:
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            self._add_title_slide(prs, title)

            total = len(slides_content)
            for idx, slide_data in enumerate(slides_content):
                self._add_content_slide(
                    prs,
                    slide_data.get("title", ""),
                    slide_data.get("content", []),
                    slide_num=idx + 1,
                    total_slides=total,
                )

            self._add_closing_slide(prs, title)

            prs.save(output_path)
            logger.info(f"Presentation saved to: {output_path}")
            return output_path

        except Exception as e:
            logger.error(f"Error creating presentation: {e}")
            raise

    # ------------------------------------------------------------------ #
    #  Title slide                                                        #
    # ------------------------------------------------------------------ #
    def _add_title_slide(self, prs, title: str):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = self.theme["title_bg"]

        # Left accent bar
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(0.35), prs.slide_height,
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.theme["accent"]
        bar.line.fill.background()

        # Title
        tb = slide.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(10), Inches(2.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = self.theme["title_color"]
        p.font.name = self.theme["font_heading"]

        # Subtitle line
        sub = slide.shapes.add_textbox(Inches(1.2), Inches(4.8), Inches(10), Inches(0.8))
        sf = sub.text_frame
        sp = sf.paragraphs[0]
        sp.text = f"Generated by CyberBron  |  {datetime.now().strftime('%B %d, %Y')}"
        sp.font.size = Pt(16)
        sp.font.color.rgb = self.theme["subtitle_color"]
        sp.font.name = self.theme["font_body"]

        # Thin accent line under title
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1.2), Inches(4.55),
            Inches(3), Pt(3),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.theme["accent"]
        line.line.fill.background()

    # ------------------------------------------------------------------ #
    #  Content slide                                                      #
    # ------------------------------------------------------------------ #
    def _add_content_slide(self, prs, title: str, content, slide_num: int, total_slides: int):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = self.theme["bg_color"]

        # Left accent bar
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(0.15), prs.slide_height,
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.theme["accent"]
        bar.line.fill.background()

        # Title background strip
        title_strip = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.15), Inches(0),
            Inches(13.183), Inches(1.4),
        )
        title_strip.fill.solid()
        title_strip.fill.fore_color.rgb = self.theme["accent_light"]
        title_strip.line.fill.background()

        # Slide title
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.25), Inches(11), Inches(1.0))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(30)
        p.font.bold = True
        p.font.color.rgb = self.theme["slide_title_color"]
        p.font.name = self.theme["font_heading"]

        # Content area
        content_left = Inches(0.8)
        content_top = Inches(1.8)
        content_width = Inches(11.5)
        content_height = Inches(5.0)

        cb = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
        cf = cb.text_frame
        cf.word_wrap = True

        if isinstance(content, str):
            self._add_paragraph_content(cf, content)
        elif isinstance(content, list):
            self._add_bullet_content(cf, content)

        # Slide number in bottom-right
        num_box = slide.shapes.add_textbox(
            Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.4)
        )
        nf = num_box.text_frame
        np_ = nf.paragraphs[0]
        np_.text = f"{slide_num} / {total_slides}"
        np_.alignment = PP_ALIGN.RIGHT
        np_.font.size = Pt(10)
        np_.font.color.rgb = self.theme["subtitle_color"]
        np_.font.name = self.theme["font_body"]

    def _add_paragraph_content(self, text_frame, content: str):
        """Add paragraph content with proper formatting and spacing."""
        paragraphs = [p.strip() for p in content.split("\n") if p.strip()]
        if not paragraphs:
            paragraphs = [content]

        first = True
        for para_text in paragraphs:
            if first:
                p = text_frame.paragraphs[0]
                first = False
            else:
                p = text_frame.add_paragraph()
            p.text = para_text
            p.font.size = Pt(16)
            p.font.color.rgb = self.theme["text_color"]
            p.font.name = self.theme["font_body"]
            p.line_spacing = Pt(24)
            p.space_after = Pt(14)

    def _add_bullet_content(self, text_frame, items: list):
        """Add bullet-point content with styled bullets."""
        first = True
        for item in items:
            if first:
                p = text_frame.paragraphs[0]
                first = False
            else:
                p = text_frame.add_paragraph()

            # Use a unicode bullet for a cleaner look
            p.text = item
            p.font.size = Pt(16)
            p.font.color.rgb = self.theme["text_color"]
            p.font.name = self.theme["font_body"]
            p.line_spacing = Pt(26)
            p.space_after = Pt(8)
            p.level = 0

            # Indent and bullet settings
            p.space_before = Pt(4)
            pPr = p._pPr
            if pPr is None:
                from pptx.oxml.ns import qn
                from lxml import etree
                pPr = etree.SubElement(p._p, qn("a:pPr"))
            pPr.set("marL", str(Emu(Inches(0.4))))
            pPr.set("indent", str(Emu(-Inches(0.25))))

            # Add bullet character via XML
            from pptx.oxml.ns import qn
            from lxml import etree
            buChar = etree.SubElement(pPr, qn("a:buChar"))
            buChar.set("char", "\u25CF")  # filled circle

            buClr = etree.SubElement(pPr, qn("a:buClr"))
            srgb = etree.SubElement(buClr, qn("a:srgbClr"))
            bc = self.theme["bullet_color"]
            srgb.set("val", f"{bc[0]:02X}{bc[1]:02X}{bc[2]:02X}")

            buSzPct = etree.SubElement(pPr, qn("a:buSzPct"))
            buSzPct.set("val", "60000")  # 60% of text size

    # ------------------------------------------------------------------ #
    #  Closing slide                                                      #
    # ------------------------------------------------------------------ #
    def _add_closing_slide(self, prs, title: str):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = self.theme["title_bg"]

        # Left accent bar
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(0.35), prs.slide_height,
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.theme["accent"]
        bar.line.fill.background()

        # "Thank You" text
        tb = slide.shapes.add_textbox(Inches(1.2), Inches(2.5), Inches(10), Inches(2))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = "Thank You"
        p.alignment = PP_ALIGN.LEFT
        p.font.size = Pt(52)
        p.font.bold = True
        p.font.color.rgb = self.theme["title_color"]
        p.font.name = self.theme["font_heading"]

        # Accent line
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1.2), Inches(4.4),
            Inches(3), Pt(3),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.theme["accent"]
        line.line.fill.background()

        # Footer
        fb = slide.shapes.add_textbox(Inches(1.2), Inches(4.8), Inches(10), Inches(0.6))
        ff = fb.text_frame
        fp = ff.paragraphs[0]
        fp.text = "Powered by CyberBron"
        fp.font.size = Pt(16)
        fp.font.color.rgb = self.theme["subtitle_color"]
        fp.font.name = self.theme["font_body"]
